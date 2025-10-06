#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint解析・JSON変換スクリプト（無料版）
python-pptxを使用してPowerPointファイルからテキストを抽出し、
Google Apps Scriptで処理しやすいJSON形式に変換する
"""

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx is not installed")
    print("Please install it with: pip install python-pptx")
    exit(1)


class PowerPointProcessor:
    """PowerPoint解析・JSON変換クラス"""

    def __init__(self):
        """初期化"""
        self.patterns = {
            # 価格パターン（強化版）
            'price': [
                r'(?:単価|価格|費用|金額)[\s:：]*[¥\\]?\s*(\d{1,3}(?:,\d{3})+)(?:\s*円)?',  # 1,000円形式
                r'[¥\\]\s*(\d{1,3}(?:,\d{3})+)',  # ¥1,000形式
                r'(\d{1,3}(?:,\d{3})+)\s*円',  # 1,000円形式
                r'(?:単価|価格|費用|金額)[\s:：]*(\d+)\s*円',  # 単価500円形式
            ],
            # 数量パターン（強化版）
            'quantity': [
                r'(?:数量|個数|枚数|ロット|部数)[\s:：]*(\d{1,3}(?:,\d{3})*)\s*(?:個|枚|部|ロット)?',
                r'(\d{1,3}(?:,\d{3})*)\s*(?:個|枚|部|ロット)',
                r'最大\s*(\d{1,3}(?:,\d{3})*)',
            ],
            # 納期パターン（強化版）
            'deadline': [
                r'(?:納期|納品|お届け)[\s:：]*(\d+)\s*(?:日|営業日|週間|ヶ月)',
                r'(\d{4}年\d{1,2}月(?:\d{1,2}日)?)\s*(?:納品|納期|想定)',
                r'(\d{1,2}月(?:上旬|中旬|下旬))',
            ],
            # 会社名パターン（強化版）
            'company': [
                r'(?:株式会社|有限会社)\s*([^\s、。\n]+)',
                r'([^\s、。\n]+)\s*(?:株式会社|有限会社)',
                r'\(株\)\s*([^\s、。\n]+)',
                r'([^\s、。\n]+)\s*\(株\)',
                r'([^\s、。\n]+様)',  # 様付き
                r'クライアント[\s:：]*([^\s、。\n]+)',
            ],
            # 日付パターン（強化版）
            'date': [
                r'(\d{4})[年/\-](\d{1,2})[月/\-](\d{1,2})日?',
                r'(\d{4})[年/](\d{1,2})月',
                r'(\d{1,2})月(\d{1,2})日',
            ],
            # イベント種別
            'event_type': [
                r'(キャンペーン|イベント|展示会|セミナー|プロモーション|運営マニュアル|進行台本|提案書|企画書)',
            ],
            # クライアント名（明示的）
            'client': [
                r'クライアント[\s:：]+([^\s\n]+)',
                r'【([^】]+)様?】',
            ],
            # 景品・ノベルティ
            'novelty': [
                r'(ノベルティ|景品|グッズ|記念品|プレゼント)',
                r'(オリジナル[^\s、。\n]{2,10})',
            ],
        }

    def extract_text_from_slide(self, slide) -> List[str]:
        """スライドからテキストを抽出"""
        texts = []

        for shape in slide.shapes:
            try:
                text = self._extract_text_from_shape(shape)
                if text.strip():
                    texts.append(text.strip())
            except Exception as e:
                # エラーが発生してもスキップして続行
                continue

        return texts

    def _extract_text_from_shape(self, shape) -> str:
        """図形からテキストを抽出"""
        text = ""

        # テキストフレーム
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text += shape.text_frame.text + "\n"

        # テーブル
        if hasattr(shape, 'table'):
            try:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            except:
                pass  # テーブルが存在しない場合はスキップ

        # グループ化された図形
        if hasattr(shape, 'shapes'):
            for sub_shape in shape.shapes:
                text += self._extract_text_from_shape(sub_shape)

        return text

    def analyze_text(self, text: str) -> Dict[str, Any]:
        """テキストから情報を抽出（強化版）"""
        info = {
            'prices': [],
            'quantities': [],
            'deadlines': [],
            'companies': [],
            'dates': [],
            'event_types': [],
            'clients': [],
            'novelties': [],
            'keywords': []
        }

        # パターンマッチング（複数パターン対応）
        for key, patterns in self.patterns.items():
            all_matches = []

            # パターンが配列の場合は全てを試す
            if isinstance(patterns, list):
                for pattern in patterns:
                    try:
                        matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
                        if matches:
                            all_matches.extend(matches)
                    except:
                        continue
            else:
                try:
                    matches = re.findall(patterns, text, re.IGNORECASE | re.MULTILINE)
                    if matches:
                        all_matches.extend(matches)
                except:
                    pass

            # データ型に応じて処理
            if key == 'price':
                prices = [self._clean_number(m) for m in all_matches]
                info['prices'] = [p for p in prices if p and p > 0]
            elif key == 'quantity':
                quantities = [self._clean_number(m) for m in all_matches]
                info['quantities'] = [q for q in quantities if q and q > 0]
            elif key == 'deadline':
                info['deadlines'] = list(set([self._clean_string(m) for m in all_matches]))
            elif key == 'company':
                info['companies'] = list(set([self._clean_string(m) for m in all_matches]))
            elif key == 'date':
                info['dates'] = self._format_dates(all_matches)
            elif key == 'event_type':
                info['event_types'] = list(set([self._clean_string(m) for m in all_matches]))
            elif key == 'client':
                info['clients'] = list(set([self._clean_string(m) for m in all_matches]))
            elif key == 'novelty':
                info['novelties'] = list(set([self._clean_string(m) for m in all_matches]))

        # キーワード抽出
        keywords = self._extract_keywords(text)
        info['keywords'] = keywords

        return info

    def _clean_number(self, number_str) -> Optional[int]:
        """数値文字列をクリーンアップ"""
        try:
            # タプルの場合は最初の要素を使用
            if isinstance(number_str, tuple):
                number_str = number_str[0] if number_str else ''

            # 文字列に変換
            number_str = str(number_str)

            # カンマ・スペースを除去して数値に変換
            cleaned = re.sub(r'[,\s]', '', number_str)
            result = int(cleaned)

            # 妥当な範囲のみ返す（1円～10億円）
            if 1 <= result <= 1000000000:
                return result
            return None
        except (ValueError, AttributeError, TypeError):
            return None

    def _clean_string(self, text) -> str:
        """文字列をクリーンアップ"""
        try:
            # タプルの場合は最初の要素を使用
            if isinstance(text, tuple):
                text = text[0] if text else ''

            # 文字列に変換して前後の空白を削除
            text = str(text).strip()

            # 不要な記号を削除
            text = re.sub(r'[\[\]（）()【】]', '', text)

            return text
        except:
            return ''

    def _format_dates(self, date_matches: List) -> List[str]:
        """日付をフォーマット"""
        formatted_dates = []
        for match in date_matches:
            try:
                if isinstance(match, tuple):
                    # (2024, 7, 15) のような形式
                    if len(match) >= 3:
                        year, month, day = match[0], match[1], match[2]
                        formatted_dates.append(f"{year}/{month.zfill(2)}/{day.zfill(2)}")
                    elif len(match) == 2:
                        year, month = match[0], match[1]
                        formatted_dates.append(f"{year}/{month.zfill(2)}/01")
                else:
                    formatted_dates.append(str(match))
            except:
                continue
        return list(set(formatted_dates))

    def _extract_keywords(self, text: str) -> List[str]:
        """キーワード抽出（簡易版）"""
        # 重要そうなキーワードを抽出
        important_words = [
            'ノベルティ', '景品', 'グッズ', 'プレゼント', 'キャンペーン',
            '展示会', 'イベント', 'セミナー', 'プロモーション',
            'エコ', '環境', 'SDGs', 'オリジナル', 'カスタム'
        ]

        keywords = []
        for word in important_words:
            if word in text:
                keywords.append(word)

        return keywords

    def process_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """PowerPointファイルを処理してJSON化"""
        try:
            presentation = Presentation(file_path)

            result = {
                'file_info': {
                    'file_name': Path(file_path).name,
                    'processed_at': datetime.now().isoformat(),
                    'slide_count': len(presentation.slides)
                },
                'slides': [],
                'summary': {
                    'all_prices': [],
                    'all_quantities': [],
                    'all_companies': [],
                    'all_keywords': [],
                    'all_deadlines': [],
                    'all_dates': [],
                    'all_event_types': [],
                    'all_clients': [],
                    'all_novelties': []
                }
            }

            # 各スライドを処理
            for i, slide in enumerate(presentation.slides, 1):
                slide_texts = self.extract_text_from_slide(slide)
                combined_text = "\n".join(slide_texts)

                analyzed_info = self.analyze_text(combined_text)

                slide_data = {
                    'slide_number': i,
                    'raw_texts': slide_texts,
                    'analyzed_info': analyzed_info,
                    'text_length': len(combined_text)
                }

                result['slides'].append(slide_data)

                # サマリー情報を蓄積
                result['summary']['all_prices'].extend(analyzed_info['prices'])
                result['summary']['all_quantities'].extend(analyzed_info['quantities'])
                result['summary']['all_companies'].extend(analyzed_info['companies'])
                result['summary']['all_keywords'].extend(analyzed_info['keywords'])
                result['summary']['all_deadlines'].extend(analyzed_info['deadlines'])
                result['summary']['all_dates'].extend(analyzed_info['dates'])
                result['summary']['all_event_types'].extend(analyzed_info['event_types'])
                result['summary']['all_clients'].extend(analyzed_info['clients'])
                result['summary']['all_novelties'].extend(analyzed_info['novelties'])

            # 重複を除去
            result['summary']['all_companies'] = list(set(result['summary']['all_companies']))
            result['summary']['all_keywords'] = list(set(result['summary']['all_keywords']))
            result['summary']['all_deadlines'] = list(set(result['summary']['all_deadlines']))
            result['summary']['all_dates'] = list(set(result['summary']['all_dates']))
            result['summary']['all_event_types'] = list(set(result['summary']['all_event_types']))
            result['summary']['all_clients'] = list(set(result['summary']['all_clients']))
            result['summary']['all_novelties'] = list(set(result['summary']['all_novelties']))

            return result

        except Exception as e:
            return {
                'error': str(e),
                'file_name': Path(file_path).name,
                'processed_at': datetime.now().isoformat()
            }


def main():
    """メイン処理"""
    processor = PowerPointProcessor()

    # 使用例
    print("PowerPoint Analysis System (Free Version)")
    print("=" * 50)

    # ファイルパスを指定（例）
    file_path = input("PowerPoint file path: ")

    if not Path(file_path).exists():
        print(f"ERROR: File not found: {file_path}")
        return

    print(f"Processing: {file_path}")
    result = processor.process_powerpoint(file_path)

    if 'error' in result:
        print(f"ERROR: {result['error']}")
        return

    # JSON出力
    output_path = Path(file_path).with_suffix('.json')
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f"Processing completed!")
    print(f"Output file: {output_path}")
    print(f"Slides: {result['file_info']['slide_count']}")
    print(f"Prices found: {len(result['summary']['all_prices'])}")
    print(f"Companies found: {len(result['summary']['all_companies'])}")


if __name__ == "__main__":
    main()