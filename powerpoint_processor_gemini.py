#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint解析・Gemini AI分析スクリプト（無料版 v4.0）
python-pptx + Gemini APIを使用してPowerPointファイルから高精度でデータを抽出
"""

import json
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx is not installed")
    print("Please install it with: pip install python-pptx")
    exit(1)

try:
    import google.generativeai as genai
except ImportError:
    print("ERROR: google-generativeai is not installed")
    print("Please install it with: pip install google-generativeai")
    exit(1)


class GeminiPowerPointProcessor:
    """Gemini API統合PowerPoint解析クラス"""

    # 無料枠の制限（Flash-Lite）
    FREE_TIER_LIMITS = {
        'daily_requests': 1000,      # 1日1,000回
        'monthly_requests': 30000,    # 月間30,000回
        'rpm': 15                     # 1分間15回
    }

    def __init__(self, api_key: Optional[str] = None, usage_log_path: Optional[str] = None):
        """
        初期化

        Args:
            api_key: Gemini APIキー（省略時は環境変数から取得）
            usage_log_path: 使用状況ログファイルパス（省略時は.gemini_usage.json）
        """
        # APIキーの設定
        self.api_key = api_key or os.environ.get('GEMINI_API_KEY')
        if not self.api_key:
            raise ValueError(
                "Gemini API key is required. "
                "Set GEMINI_API_KEY environment variable or pass api_key parameter."
            )

        # 使用状況ログの設定
        self.usage_log_path = usage_log_path or Path.home() / '.gemini_usage.json'
        self.usage_data = self._load_usage_data()

        # Gemini APIの初期化
        genai.configure(api_key=self.api_key)
        # 無料版推奨モデル: Flash-Lite (1日1,000回、月30,000回まで)
        self.model = genai.GenerativeModel('gemini-2.0-flash-lite')

        print("Gemini API initialized successfully (using gemini-2.0-flash-lite)")
        self._print_usage_status()

    def _load_usage_data(self) -> Dict[str, Any]:
        """使用状況データを読み込み"""
        if Path(self.usage_log_path).exists():
            try:
                with open(self.usage_log_path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except:
                pass

        # 初期データ
        return {
            'daily': {},
            'monthly': {},
            'total': 0
        }

    def _save_usage_data(self):
        """使用状況データを保存"""
        try:
            with open(self.usage_log_path, 'w', encoding='utf-8') as f:
                json.dump(self.usage_data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"WARNING: Failed to save usage data: {e}")

    def _check_free_tier_limit(self) -> bool:
        """無料枠の制限チェック（超過したらFalseを返す）"""
        today = datetime.now().strftime('%Y-%m-%d')
        this_month = datetime.now().strftime('%Y-%m')

        # 日次使用量
        daily_count = self.usage_data['daily'].get(today, 0)
        if daily_count >= self.FREE_TIER_LIMITS['daily_requests']:
            print(f"\n⚠️  FREE TIER LIMIT EXCEEDED: Daily limit ({self.FREE_TIER_LIMITS['daily_requests']} requests/day)")
            print(f"   Today's usage: {daily_count}/{self.FREE_TIER_LIMITS['daily_requests']}")
            print(f"   システムを停止します（無料枠超過のため課金を防止）")
            return False

        # 月次使用量
        monthly_count = self.usage_data['monthly'].get(this_month, 0)
        if monthly_count >= self.FREE_TIER_LIMITS['monthly_requests']:
            print(f"\n⚠️  FREE TIER LIMIT EXCEEDED: Monthly limit ({self.FREE_TIER_LIMITS['monthly_requests']} requests/month)")
            print(f"   This month's usage: {monthly_count}/{self.FREE_TIER_LIMITS['monthly_requests']}")
            print(f"   システムを停止します（無料枠超過のため課金を防止）")
            return False

        return True

    def _increment_usage(self):
        """使用回数をカウント"""
        today = datetime.now().strftime('%Y-%m-%d')
        this_month = datetime.now().strftime('%Y-%m')

        # 日次カウント
        self.usage_data['daily'][today] = self.usage_data['daily'].get(today, 0) + 1

        # 月次カウント
        self.usage_data['monthly'][this_month] = self.usage_data['monthly'].get(this_month, 0) + 1

        # 総カウント
        self.usage_data['total'] = self.usage_data.get('total', 0) + 1

        # 保存
        self._save_usage_data()

    def _print_usage_status(self):
        """現在の使用状況を表示"""
        today = datetime.now().strftime('%Y-%m-%d')
        this_month = datetime.now().strftime('%Y-%m')

        daily_count = self.usage_data['daily'].get(today, 0)
        monthly_count = self.usage_data['monthly'].get(this_month, 0)
        total_count = self.usage_data.get('total', 0)

        daily_remaining = self.FREE_TIER_LIMITS['daily_requests'] - daily_count
        monthly_remaining = self.FREE_TIER_LIMITS['monthly_requests'] - monthly_count

        print(f"\n📊 FREE TIER USAGE STATUS:")
        print(f"   Today: {daily_count}/{self.FREE_TIER_LIMITS['daily_requests']} requests (残り {daily_remaining})")
        print(f"   This month: {monthly_count}/{self.FREE_TIER_LIMITS['monthly_requests']} requests (残り {monthly_remaining})")
        print(f"   Total: {total_count} requests")
        print()

    def extract_text_from_slide(self, slide) -> List[str]:
        """スライドからテキストを抽出"""
        texts = []

        for shape in slide.shapes:
            try:
                text = self._extract_text_from_shape(shape)
                if text.strip():
                    texts.append(text.strip())
            except Exception as e:
                continue

        return texts

    def _extract_text_from_shape(self, shape) -> str:
        """図形からテキストを抽出"""
        text = ""

        # テキストフレーム
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text += shape.text_frame.text + "\n"

        # テーブル
        if hasattr(shape, 'table'):
            try:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            except:
                pass

        # グループ化された図形
        if hasattr(shape, 'shapes'):
            for sub_shape in shape.shapes:
                text += self._extract_text_from_shape(sub_shape)

        return text

    def analyze_with_gemini(self, slide_texts: List[str], file_name: str) -> Dict[str, Any]:
        """
        Gemini APIでテキストを分析

        Args:
            slide_texts: スライドのテキストリスト
            file_name: ファイル名

        Returns:
            構造化された分析結果
        """
        # 無料枠チェック
        if not self._check_free_tier_limit():
            error_result = self._get_empty_analysis()
            error_result['error'] = 'FREE_TIER_LIMIT_EXCEEDED'
            return error_result

        # テキストを結合
        combined_text = "\n\n".join(slide_texts)

        # ファイル名からクライアント名を事前抽出
        client_hint = self._extract_client_from_filename(file_name)

        # プロンプト作成
        prompt = f"""あなたはプロモーション事業のデータ分析AIです。
以下のPowerPointスライドのテキストから、構造化データを抽出してください。

【ファイル名】
{file_name}
{f'（クライアント名ヒント: {client_hint}）' if client_hint else ''}

【抽出項目】
1. client_name: クライアント名（【XX様】などから企業名を抽出。「様」「株式会社」「有限会社」は除く）
2. event_date: 実施時期（YYYY/MM/DD形式で。複数ある場合は最も重要なもの）
3. event_type: イベント種別（提案書/運営マニュアル/進行台本/企画書/キャンペーン/イベント/展示会/セミナーなど）
4. event_description: イベント内容の概要（1-2文で）
5. unit_price: 単価（円、数値のみ。複数ある場合は代表的なもの）
6. total_cost: 総費用（円、数値のみ）
7. order_quantity: 発注数量（数値のみ）
8. target_count: ターゲット人数（「先着XX名」などから）
9. deadline: 納期（「XX営業日」「YYYY年MM月」など、元の表現を保持）
10. partner_companies: 協力会社名のリスト（最大5社）
11. novelty_items: ノベルティ/景品の具体的な名称リスト（最大5個）
12. venue: 会場名
13. keywords: 重要なキーワードリスト（最大10個）

【スライドテキスト】
{combined_text[:3000]}

【出力形式】
以下のJSON形式で出力してください。値が不明な場合はnullを設定してください。
{{
  "client_name": "クライアント名",
  "event_date": "2024/01/01",
  "event_type": "種別",
  "event_description": "概要",
  "unit_price": 500,
  "total_cost": 300000,
  "order_quantity": 1000,
  "target_count": 500,
  "deadline": "14営業日",
  "partner_companies": ["会社1", "会社2"],
  "novelty_items": ["景品1", "景品2"],
  "venue": "会場名",
  "keywords": ["キーワード1", "キーワード2"]
}}

重要: 必ずJSON形式のみを出力してください。説明文は不要です。"""

        try:
            # Gemini APIに送信
            print("  Sending to Gemini API...")
            response = self.model.generate_content(prompt)

            # API使用回数をカウント（成功したらカウント）
            self._increment_usage()

            # レスポンスからJSONを抽出
            response_text = response.text.strip()

            # JSONブロックを抽出（```json ``` で囲まれている場合）
            json_match = re.search(r'```json\s*(.*?)\s*```', response_text, re.DOTALL)
            if json_match:
                json_str = json_match.group(1)
            else:
                # JSONブロックがない場合はそのまま使用
                json_str = response_text

            # JSONをパース
            analyzed_data = json.loads(json_str)

            # 信頼度スコアを計算
            confidence = self._calculate_confidence(analyzed_data)
            analyzed_data['confidence_score'] = confidence

            print(f"  Gemini API analysis completed (confidence: {confidence}%)")
            return analyzed_data

        except json.JSONDecodeError as e:
            print(f"  ERROR: Failed to parse Gemini response as JSON: {e}")
            print(f"  Response: {response_text[:500]}")
            return self._get_empty_analysis()
        except Exception as e:
            print(f"  ERROR: Gemini API call failed: {e}")
            return self._get_empty_analysis()

    def _extract_client_from_filename(self, filename: str) -> str:
        """ファイル名からクライアント名を抽出"""
        # 【クライアント名様】パターン
        match = re.search(r'【([^】]+)様?】', filename)
        if match:
            return match.group(1)
        return ""

    def _calculate_confidence(self, data: Dict[str, Any]) -> int:
        """信頼度スコアを計算"""
        score = 0

        # 各項目の有無をチェック
        if data.get('client_name'): score += 15
        if data.get('event_date'): score += 15
        if data.get('event_type'): score += 10
        if data.get('event_description'): score += 10
        if data.get('unit_price'): score += 10
        if data.get('total_cost'): score += 10
        if data.get('order_quantity'): score += 5
        if data.get('deadline'): score += 5
        if data.get('partner_companies') and len(data['partner_companies']) > 0: score += 10
        if data.get('novelty_items') and len(data['novelty_items']) > 0: score += 5
        if data.get('keywords') and len(data['keywords']) > 0: score += 5

        return min(score, 100)

    def _get_empty_analysis(self) -> Dict[str, Any]:
        """空の分析結果を返す"""
        return {
            'client_name': None,
            'event_date': None,
            'event_type': None,
            'event_description': None,
            'unit_price': None,
            'total_cost': None,
            'order_quantity': None,
            'target_count': None,
            'deadline': None,
            'partner_companies': [],
            'novelty_items': [],
            'venue': None,
            'keywords': [],
            'confidence_score': 0
        }

    def process_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """PowerPointファイルを処理してJSON化"""
        try:
            print(f"Processing: {Path(file_path).name}")
            presentation = Presentation(file_path)

            # 全スライドからテキストを抽出
            all_slide_texts = []
            for i, slide in enumerate(presentation.slides, 1):
                slide_texts = self.extract_text_from_slide(slide)
                all_slide_texts.extend(slide_texts)
                print(f"  Slide {i}/{len(presentation.slides)}: {len(slide_texts)} text blocks extracted")

            # Gemini APIで分析
            analyzed_data = self.analyze_with_gemini(all_slide_texts, Path(file_path).name)

            # 結果を構築
            result = {
                'file_info': {
                    'file_name': Path(file_path).name,
                    'processed_at': datetime.now().isoformat(),
                    'slide_count': len(presentation.slides),
                    'processing_method': 'gemini_api_v4.0'
                },
                'gemini_analysis': analyzed_data,
                'slide_texts_sample': '\n'.join(all_slide_texts[:5])[:1000]  # サンプルのみ保存
            }

            return result

        except Exception as e:
            return {
                'error': str(e),
                'file_name': Path(file_path).name,
                'processed_at': datetime.now().isoformat()
            }


def main():
    """メイン処理"""
    print("=" * 60)
    print("PowerPoint Analysis System (Gemini API v4.0)")
    print("=" * 60)

    # APIキーの確認
    api_key = os.environ.get('GEMINI_API_KEY')
    if not api_key:
        print("\nGemini API key not found in environment variable.")
        api_key = input("Enter your Gemini API key: ").strip()
        if not api_key:
            print("ERROR: API key is required")
            return

    # プロセッサー初期化
    try:
        processor = GeminiPowerPointProcessor(api_key=api_key)
    except Exception as e:
        print(f"ERROR: Failed to initialize: {e}")
        return

    # ファイルパスを入力
    file_path = input("\nPowerPoint file path: ").strip()

    if not Path(file_path).exists():
        print(f"ERROR: File not found: {file_path}")
        return

    # 処理実行
    result = processor.process_powerpoint(file_path)

    if 'error' in result:
        print(f"\nERROR: {result['error']}")
        return

    # JSON出力
    output_path = Path(file_path).with_suffix('.json')
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    # 結果表示
    print("\n" + "=" * 60)
    print("Processing completed!")
    print("=" * 60)
    print(f"Output file: {output_path}")
    print(f"Slides: {result['file_info']['slide_count']}")

    analysis = result['gemini_analysis']
    print(f"\n【AI Analysis Results】")
    print(f"Client: {analysis.get('client_name', 'N/A')}")
    print(f"Event Date: {analysis.get('event_date', 'N/A')}")
    print(f"Event Type: {analysis.get('event_type', 'N/A')}")
    print(f"Unit Price: {analysis.get('unit_price', 'N/A')} yen")
    print(f"Partner Companies: {len(analysis.get('partner_companies', []))} found")
    print(f"Confidence Score: {analysis.get('confidence_score', 0)}%")


if __name__ == "__main__":
    main()